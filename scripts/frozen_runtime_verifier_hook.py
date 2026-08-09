"""Opt-in PyInstaller runtime hooks for release artifact verification.

Both modes below are inert for ordinary application launches.

``MATRX_FROZEN_RUNTIME_VERIFY=1`` — release/smoke tooling points
``MATRX_FROZEN_RUNTIME_PATH`` at the exact locked image-generation environment.
The directory is appended, matching production precedence, then critical lazy
imports execute inside the real frozen CPython process.

``MATRX_FROZEN_OFFICE_VERIFY=1`` — exercises the Office (docx/pptx/xlsx) codec
end to end inside the frozen process: generate from a spec, author a document
with the raw renderers, then extract it back to markdown. Deliberately
independent of the managed image runtime so it gates EVERY target, including
the ones whose managed runtime is unsupported and archive-only.
"""

from __future__ import annotations

import importlib
import importlib.metadata
import hashlib
import json
import os
import sys
import traceback
from pathlib import Path


def _verify_office_codec() -> dict[str, object]:
    """Prove the frozen bundle can read AND write Office documents.

    Every module here is lazily imported inside a function in production
    (``app.tools.tools.file_ops._read_office`` and
    ``app.tools.tools.media.tool_office_generate``), so PyInstaller's static
    analysis reaches none of it. The specs compensate with explicit
    ``collect_submodules``/``collect_data_files``; this probe is what proves
    they actually worked in the artifact users receive.

    ``docx.Document()`` and ``pptx.Presentation()`` with no argument load
    ``docx/templates/default.docx`` / ``pptx/templates/default.pptx`` from
    beside the package — data files, not modules — so a bundle can carry every
    module and still fail here.
    """
    import io
    import mimetypes

    from matrx_files.specific_handlers.office import (
        DocumentSpec,
        PresentationSpec,
        SpreadsheetSpec,
        classify_office,
        extract_office,
        generate_office,
    )

    details: dict[str, object] = {}

    # ── Routing: the decision tool_read makes before the codec is reached ────
    # A frozen process may resolve `mimetypes` differently from the dev host
    # (no /etc/mime.types, different registry). If classification returns None,
    # tool_read falls through to the bounded text reader and hands the model
    # zip garbage instead of markdown — a silent wrong answer, not an error.
    routing: dict[str, str] = {}
    for extension in (".docx", ".pptx", ".xlsx"):
        guessed, _ = mimetypes.guess_type(f"document{extension}")
        kind = classify_office(guessed, f"document{extension}")
        if kind is None:
            raise RuntimeError(
                f"classify_office did not recognize {extension} "
                f"(mimetypes guessed {guessed!r}) — tool_read would read it as text"
            )
        routing[extension] = kind.value
    details["routing"] = routing

    # ── The default templates must exist as bundled data files ───────────────
    import docx as _docx
    import pptx as _pptx

    templates = {
        "docx": Path(_docx.__file__).parent / "templates" / "default.docx",
        "pptx": Path(_pptx.__file__).parent / "templates" / "default.pptx",
    }
    for kind, template in templates.items():
        if not template.is_file():
            raise RuntimeError(
                f"bundled {kind} default template missing: {template} — "
                f"collect_data_files({kind!r}) did not reach the artifact"
            )
    details["templates"] = {k: str(v) for k, v in templates.items()}

    # ── generate: the tool_office_generate path, all three formats ───────────
    generated: dict[str, int] = {}
    specs = {
        "docx": DocumentSpec.model_validate(
            {
                "title": "Frozen Bundle Probe",
                "blocks": [
                    {"type": "heading", "text": "Office Codec", "level": 1},
                    {"type": "paragraph", "text": "generated inside the frozen engine"},
                    {
                        "type": "table",
                        "header": True,
                        "rows": [
                            ["Format", "Status"],
                            ["docx", "ok"],
                            ["pptx", "ok"],
                            ["xlsx", "ok"],
                        ],
                    },
                ],
            }
        ),
        "pptx": PresentationSpec.model_validate(
            {
                "title": "Frozen Bundle Probe",
                "subtitle": "office codec",
                "slides": [
                    {
                        "title": "Office Codec",
                        "bullets": ["generated inside the frozen engine"],
                        "notes": "probe slide",
                    }
                ],
            }
        ),
        "xlsx": SpreadsheetSpec.model_validate(
            {
                "sheets": [
                    {
                        "name": "Office Codec",
                        "columns": ["Format", "Status"],
                        "rows": [["docx", "ok"], ["pptx", "ok"], ["xlsx", "ok"]],
                    }
                ]
            }
        ),
    }
    import base64

    for kind, spec in specs.items():
        result = generate_office(spec)
        payload = base64.b64decode(result.base64_data)
        if not payload.startswith(b"PK"):
            raise RuntimeError(f"generated {kind} is not an OpenXML package")
        generated[kind] = len(payload)
        # Round-trip the freshly generated bytes back through the reader.
        round_trip = extract_office(payload, file_name=f"probe.{kind}")
        if "Office Codec" not in (round_trip.markdown or ""):
            raise RuntimeError(
                f"generated {kind} did not survive extraction: "
                f"{(round_trip.markdown or '')[:200]!r}"
            )
    details["generated_bytes"] = generated

    # ── extract: documents authored with the raw renderers, carrying shapes
    # the generator never emits (heading levels, speaker notes, a second sheet,
    # a formula) so extraction is not merely reading back its own output ──────
    document = _docx.Document()
    # Headers/footers load docx/parts/../templates/default-{header,footer}.xml —
    # a path that walks UP out of a subpackage, and therefore only resolves when
    # docx/parts/ physically exists under sys._MEIPASS. Same shape as the pptx
    # notes lookup above; both are frozen-only failures.
    section = document.sections[0]
    section.header.is_linked_to_previous = False
    section.header.paragraphs[0].text = "Probe header"
    section.footer.is_linked_to_previous = False
    section.footer.paragraphs[0].text = "Probe footer"
    document.add_heading("Quarterly Review", level=1)
    document.add_heading("Highlights", level=2)
    document.add_paragraph("Revenue grew across every region.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Region"
    table.cell(0, 1).text = "Growth"
    table.cell(1, 0).text = "West"
    table.cell(1, 1).text = "12%"
    buffer = io.BytesIO()
    document.save(buffer)
    docx_markdown = extract_office(buffer.getvalue(), file_name="review.docx").markdown or ""
    for needle in ("Quarterly Review", "Highlights", "Revenue grew", "West"):
        if needle not in docx_markdown:
            raise RuntimeError(f"docx extraction lost {needle!r}: {docx_markdown[:400]!r}")

    presentation = _pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Roadmap"
    slide.placeholders[1].text = "Ship the Office codec"
    slide.notes_slide.notes_text_frame.text = "speaker note"
    buffer = io.BytesIO()
    presentation.save(buffer)
    pptx_markdown = extract_office(buffer.getvalue(), file_name="deck.pptx").markdown or ""
    for needle in ("Roadmap", "Ship the Office codec"):
        if needle not in pptx_markdown:
            raise RuntimeError(f"pptx extraction lost {needle!r}: {pptx_markdown[:400]!r}")

    import openpyxl

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Numbers"
    sheet.append(["Item", "Count"])
    sheet.append(["widgets", 4])
    sheet.append(["total", "=SUM(B2:B2)"])
    workbook.create_sheet("Notes").append(["second sheet"])
    buffer = io.BytesIO()
    workbook.save(buffer)
    xlsx_markdown = extract_office(buffer.getvalue(), file_name="book.xlsx").markdown or ""
    for needle in ("Numbers", "widgets", "second sheet"):
        if needle not in xlsx_markdown:
            raise RuntimeError(f"xlsx extraction lost {needle!r}: {xlsx_markdown[:400]!r}")

    details["extracted_chars"] = {
        "docx": len(docx_markdown),
        "pptx": len(pptx_markdown),
        "xlsx": len(xlsx_markdown),
    }
    # Best-effort: .dist-info is only bundled for packages that read their own
    # metadata at import time, and none of these do. A missing version here is
    # cosmetic and must never fail the gate.
    versions: dict[str, str] = {}
    for name in ("matrx-files", "python-docx", "python-pptx", "openpyxl"):
        try:
            versions[name] = importlib.metadata.version(name)
        except Exception:
            versions[name] = "unknown"
    details["versions"] = versions
    return details


if os.environ.get("MATRX_FROZEN_OFFICE_VERIFY") == "1":
    office_result: dict[str, object] = {"ok": False}
    try:
        office_result.update(_verify_office_codec())
        office_result["ok"] = True
    except BaseException as exc:  # release gate must serialize every failure
        office_result["error"] = f"{type(exc).__name__}: {exc}"
        office_result["traceback"] = traceback.format_exc()

    print(
        "MATRX_FROZEN_OFFICE_VERIFY=" + json.dumps(office_result, sort_keys=True),
        flush=True,
    )
    sys.stdout.flush()
    sys.stderr.flush()
    os._exit(0 if office_result["ok"] else 87)


if os.environ.get("MATRX_FROZEN_RUNTIME_VERIFY") == "1":
    result: dict[str, object] = {
        "contract": "unknown",
        "ok": False,
        "python": f"{sys.version_info.major}.{sys.version_info.minor}",
    }
    try:
        runtime_path = Path(os.environ["MATRX_FROZEN_RUNTIME_PATH"]).resolve(strict=True)
        if not runtime_path.is_dir():
            raise RuntimeError(f"managed runtime is not a directory: {runtime_path}")
        runtime_text = str(runtime_path)
        while runtime_text in sys.path:
            sys.path.remove(runtime_text)
        # Production contract: frozen/core dependencies win. Never prepend here.
        sys.path.append(runtime_text)

        base = Path(getattr(sys, "_MEIPASS", Path(__file__).resolve().parents[1]))
        contract_path = (
            base / "config" / "runtime-manifests" / "image-gen-contract.json"
        )
        contract = json.loads(contract_path.read_text(encoding="utf-8"))
        result["contract"] = contract["contract_sha256"]
        if result["python"] != contract["python_minor"]:
            raise RuntimeError(
                f"frozen Python {result['python']} != contract {contract['python_minor']}"
            )

        target = os.environ["MATRX_FROZEN_RUNTIME_TARGET"]
        target_path = (
            base / "config" / "runtime-manifests" / f"image-gen-{target}.json"
        )
        target_manifest = json.loads(target_path.read_text(encoding="utf-8"))
        if target_manifest.get("supported") is not True:
            raise RuntimeError(f"verification invoked for unsupported target {target}")
        if target_manifest.get("target") != target:
            raise RuntimeError("embedded target manifest identifies the wrong platform")
        if target_manifest.get("contract_sha256") != contract["contract_sha256"]:
            raise RuntimeError("embedded target manifest is stale")
        lock_path = target_path.parent / target_manifest["lock_file"]
        if hashlib.sha256(lock_path.read_bytes()).hexdigest() != target_manifest.get(
            "lock_sha256"
        ):
            raise RuntimeError("embedded target requirements digest is invalid")
        expected_versions = {
            item["name"].lower().replace("_", "-"): item["version"]
            for item in target_manifest["packages"]
        }
        actual_versions = {
            distribution.metadata["Name"].lower().replace("_", "-"): distribution.version
            for distribution in importlib.metadata.distributions(path=[runtime_text])
        }
        if actual_versions != expected_versions:
            raise RuntimeError(
                "isolated runtime package inventory differs from target manifest: "
                f"expected={expected_versions!r}, actual={actual_versions!r}"
            )
        frozen_shared_versions = contract["shared_versions_by_target"][target]
        for name, version in frozen_shared_versions.items():
            if expected_versions.get(name) != version:
                raise RuntimeError(
                    f"managed lock shared {name}={expected_versions.get(name)!r}; "
                    f"frozen contract requires {version!r}"
                )

        imported: dict[str, str] = {}
        for module_name in contract["runtime_imports"]:
            module = importlib.import_module(module_name)
            imported[module_name] = str(getattr(module, "__file__", "<built-in>"))
        for module_name, attributes in contract["runtime_attributes"].items():
            module = importlib.import_module(module_name)
            for attribute in attributes:
                getattr(module, attribute)

        frozen_shared_origins: dict[str, str] = {}
        for module_name in contract["shared_import_packages_by_target"][target]:
            module = importlib.import_module(module_name)
            origin = Path(str(getattr(module, "__file__", ""))).resolve(strict=False)
            if origin.is_relative_to(runtime_path):
                raise RuntimeError(
                    f"shared module {module_name} resolved from managed runtime: {origin}"
                )
            frozen_shared_origins[module_name] = str(origin)

        for module_name in (
            "accelerate",
            "diffusers",
            "gguf",
            "peft",
            "sentencepiece",
            "torch",
            "torchvision",
            "transformers",
        ):
            origin = Path(imported[module_name]).resolve(strict=False)
            if not origin.is_relative_to(runtime_path):
                raise RuntimeError(
                    f"managed module {module_name} resolved outside isolated runtime: {origin}"
                )

        import torch
        import torchvision

        tensor = torch.tensor([1.0, 2.0]) + 1
        if tensor.tolist() != [2.0, 3.0]:
            raise RuntimeError("PyTorch CPU operation returned an unexpected result")
        # Torchvision loads its `_C` library through torch.ops rather than as a
        # normal PyInit extension. `_has_ops()` proves that native activation.
        if not torchvision.extension._has_ops():
            raise RuntimeError("Torchvision native operators failed to load")
        torch_variant = target_manifest["torch_variant"]
        if torch_variant == "cu126":
            if "+cu126" not in str(torch.__version__) or "+cu126" not in str(
                torchvision.__version__
            ):
                raise RuntimeError("PyTorch/Torchvision are not cu126 builds")
            if not str(torch.version.cuda or "").startswith("12.6"):
                raise RuntimeError(
                    f"PyTorch reports CUDA {torch.version.cuda!r}, expected 12.6"
                )
        elif torch_variant == "mps" and not torch.backends.mps.is_built():
            raise RuntimeError("PyTorch was not built with required MPS support")

        result["imports"] = imported
        result["frozen_shared_origins"] = frozen_shared_origins
        result["torch_variant"] = torch_variant
        result["target"] = target
        result["torch"] = str(torch.__version__)
        result["torchvision"] = str(torchvision.__version__)
        result["ok"] = True
    except BaseException as exc:  # release gate must serialize every failure
        result["error"] = f"{type(exc).__name__}: {exc}"
        result["traceback"] = traceback.format_exc()

    print("MATRX_FROZEN_RUNTIME_VERIFY=" + json.dumps(result, sort_keys=True), flush=True)
    sys.stdout.flush()
    sys.stderr.flush()
    os._exit(0 if result["ok"] else 86)
